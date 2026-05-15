from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(prs, title, content, is_table=False, headers=None, rows=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)
    
    if is_table and headers:
        cols = len(headers)
        tbl = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5)).table
        for i, h in enumerate(headers):
            c = tbl.cell(0, i)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(255, 215, 0)
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = tbl.cell(ri + 1, ci)
                c.text = str(val)
                p = c.text_frame.paragraphs[0]
                p.font.size = Pt(22)
                p.alignment = PP_ALIGN.CENTER
    else:
        cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
        tf = cb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(28)
            p.space_after = Pt(20)

# 1. 표지
add_slide(prs, 'HeartChain', ['블록체인 자원봉사 플랫폼', '기여도는 영구적으로 기록', '노동이 진정한 가치를 창조', '플랫폼 성장의 수익을 공유'])

# 2. 핵심 가치 - 팀
add_slide(prs, '핵심 가치 - 팀', ['플랫폼 지분 확보', '블록체인 인증', '무료 태스크 게시', '태스크의 10% 보상'])

# 3. 핵심 가치 - 개인
add_slide(prs, '핵심 가치 - 개인', ['영구적인 기여 기록', '포인트로 상품 교환', '노동이 가치를 창조', '기술로 수익 창출'])

# 4. 지분 구조
add_slide(prs, '지분 구조', [], True, ['지분 보유자', '비율', '비고'], [['팀', '50%', '팀원 수比例'], ['개인', '20%', '균등'], ['개발', '30%', '운영']])

# 5. 지분 계산 방법
add_slide(prs, '지분 계산 방법', ['팀 지분 = 50% x 팀 인원 / 전체', '개인 지분 = 20% / 전체 인원'])

# 6. 지분 예시
add_slide(prs, '지분 예시', [], True, ['단계', '회원 수', '비율', '지분'], [['初期', '1K', '10%', '5%'], ['성장', '10K', '20%', '10%'], ['성숙', '100K', '30%', '15%']])

# 7. 팀 특권
add_slide(prs, '팀 특권', [], True, ['기능', '팀', '개인'], [['게시', '무료', '유료'], ['보상', '10%', '없음']])

# 8. 포인트 시스템
add_slide(prs, '포인트 시스템', ['포인트 = 노동 가치', '가치 = 최저임금 x 비율', '태스크 및 몰에서 사용 가능'])

# 9. 제네시스 인센티브
add_slide(prs, '제네시스 인센티브', ['최초 10개 팀 한정', '50% 추가 지분', '설립자 배지', '조기 접근 권한'])

# 10. 가치 추정
add_slide(prs, '가치 추정', [], True, ['보유자', '지분', '가치', '월간'], [['팀', '50%', '5억', '900만'], ['회원', '20%', '2억', '360만'], ['개발', '30%', '3억', '540만']])

# 11. 안전 보장
add_slide(prs, '안전 보장', ['활성 회원: 6개월内有 로그인', '비활성 계정 자동 삭제', '사진 인증 필요', '이중 서명 확인'])

# 12. 로드맵
add_slide(prs, '로드맵', [], True, ['단계', '시점', '목표'], [['제네시스', '2026 Q2', '1만'], ['선구자', '2026 Q3', '10만'], ['성장', '2026 Q4', '100만'], ['런칭', '2027 Q1', '1000만']])

# 13. 연락처
add_slide(prs, '연락처', ['모든 헌신을 기록하다', '劳动에 가치를 부여하다', '지금联系我们하세요'])

prs.save('E:/WorkBuddy/heartchain/docs/HeartChain_Korean.pptx')
print('完成！한국어 버전 PPT 생성 완료!')
