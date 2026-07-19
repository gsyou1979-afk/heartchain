from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(prs, title, content, is_table=False, headers=None, rows=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)
    
    if is_table and headers:
        cols = len(headers)
        tbl = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5)).table
        for i, h in enumerate(headers):
            c = tbl.cell(0, i)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(255, 215, 0)
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(24)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                c = tbl.cell(ri + 1, ci)
                c.text = str(val)
                p = c.text_frame.paragraphs[0]
                p.font.size = Pt(22)
                p.alignment = PP_ALIGN.CENTER
    else:
        cb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
        tf = cb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(28)
            p.space_after = Pt(20)

add_slide(prs, 'HeartChain', ['Blockchain Volunteer Platform', 'Contributions recorded permanently', 'Labor creates real value', 'Share platform growth'])
add_slide(prs, 'Core Value - Teams', ['Get platform shares', 'Blockchain endorsement', 'Free task posting', '10% reward on tasks'])
add_slide(prs, 'Core Value - Individuals', ['Permanent contribution record', 'Points redeemable', 'Labor creates value', 'Skill monetization'])
add_slide(prs, 'Equity Structure', [], True, ['Holder', 'Ratio', 'Note'], [['Teams', '50%', 'By member'], ['Individuals', '20%', 'Equal'], ['Dev', '30%', 'Ops']])
add_slide(prs, 'How to Calculate', ['Team shares = 50% x team / total', 'Individual = 20% / total'])
add_slide(prs, 'Share Example', [], True, ['Stage', 'Members', 'Ratio', 'Shares'], [['Early', '1K', '10%', '5%'], ['Growth', '10K', '20%', '10%'], ['Mature', '100K', '30%', '15%']])
add_slide(prs, 'Team Privileges', [], True, ['Feature', 'Team', 'Individual'], [['Post', 'Free', 'Pay'], ['Reward', '10%', 'None']])
add_slide(prs, 'Points System', ['Points = Labor Value', 'Value = Min wage x ratio', 'Use for tasks and mall'])
add_slide(prs, 'Genesis Incentives', ['First 10 teams', '50% extra shares', 'Founder badge', 'Early access'])
add_slide(prs, 'Value Estimate', [], True, ['Holder', 'Shares', 'Value', 'Monthly'], [['Teams', '50%', '500M', '9M'], ['Members', '20%', '200M', '3.6M'], ['Dev', '30%', '300M', '5.4M']])
add_slide(prs, 'Security', ['Active members: 6 months login', 'Auto remove inactive', 'Photo verification', 'Dual sign-off'])
add_slide(prs, 'Roadmap', [], True, ['Phase', 'Time', 'Target'], [['Genesis', 'Q2 2026', '10K'], ['Pioneer', 'Q3 2026', '100K'], ['Growth', 'Q4 2026', '1M'], ['Launch', 'Q1 2027', '10M']])
add_slide(prs, 'Contact Us', ['Record every spirit', 'Make labor valuable', 'Contact us now'])

prs.save('E:/WorkBuddy/heartchain/docs/HeartChain_Project.pptx')
print('Done!')